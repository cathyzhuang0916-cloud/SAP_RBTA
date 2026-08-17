"""
KPI calculation engine — shared by update_pptx.py and app.py.

compute_kpis(excel_path)  → dict with all KPI values + timestamp
build_pptx(kpi_data, template_path, output_path)  → writes updated PPTX
"""

import warnings
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

warnings.filterwarnings("ignore")

# ── Colors ───────────────────────────────────────────────────────────────────
GREEN = RGBColor(0x34, 0xA4, 0x1C)
AMBER = RGBColor(0xFF, 0xC0, 0x00)
RED   = RGBColor(0xFF, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)

REGIONS = ["ANZ", "G.China", "India", "Japan", "Korea", "South East Asia"]
ALL_REGIONS = REGIONS + ["APAC"]
REG_COL = "SBB_REGION_LEVEL_02"

Q3_TARGETS = {
    "joule":        "65%",
    "ai_plan":      "75%",
    "rwsm":         "55%",
    "risk":         "50%",
    "calm":         "90%",
    "gainsight":    "90%",
    "signavio":     "65%",
    "leanix":       "50%",
    "ea_cust":      "31.5%",
    "ea_acv":       "74%",
    "ea_density":   "3.2",
    "ea_acv_den":   "16.5 M€",
    "leanix_enable":"85%",
    "ea_learning":  "85%",
}

# Frozen values (Gainsight + Learning — not recalculated from data)
FROZEN = {
    "gainsight":     {"ANZ":"56%","G.China":"100%","India":"79%","Japan":"95%","Korea":"100%","South East Asia":"83%","APAC":"81%"},
    "leanix_enable": {"ANZ":"78%","G.China":"56%","India":"68%","Japan":"58%","Korea":"75%","South East Asia":"71%","APAC":"68%"},
    "ea_learning":   {"ANZ":"78%","G.China":"56%","India":"68%","Japan":"58%","Korea":"75%","South East Asia":"71%","APAC":"68%"},
}

# ── Helpers ──────────────────────────────────────────────────────────────────

def _pct(num, den):
    if den == 0:
        return None
    return round(100.0 * num / den)

def _pct_str(num, den):
    v = _pct(num, den)
    return f"{v}%" if v is not None else "N/A"

def _is_true(v):
    return str(v).strip() == "True"

def _is_yes(v):
    return str(v).strip().upper() == "YES"

def _delta(new_str, old_str):
    try:
        n = float(str(new_str).replace("%","").replace("M€","").strip())
        o = float(str(old_str).replace("%","").replace("M€","").strip())
        diff = round(n - o, 2)
        if diff == 0:
            return "(⬤)"
        sign = "+" if diff > 0 else ""
        unit = "%" if "%" in str(new_str) else ""
        return f"({sign}{int(diff) if diff == int(diff) else diff}{unit})"
    except Exception:
        return "(⬤)"

def _mom_bracket(new_str, old_str):
    """Return MoM bracket string like '+9% MoM' or '-3% MoM', or '' if zero/not computable."""
    try:
        n = float(str(new_str).replace("%","").replace("M€","").strip())
        o = float(str(old_str).replace("%","").replace("M€","").strip())
        diff = round(n - o, 2)
        if diff == 0:
            return ""
        sign = "+" if diff > 0 else ""
        unit = "%" if "%" in str(new_str) else ""
        diff_fmt = int(diff) if diff == int(diff) else diff
        return f"{sign}{diff_fmt}{unit} MoM"
    except Exception:
        return ""

def get_color_str(actual_str, target_str):
    """Return 'green', 'amber', or 'red' based on actual vs target (≥target=green, 70–99%=amber, <70%=red)."""
    try:
        def _clean(s):
            return str(s).replace("%","").replace("M€","").replace("ME","").replace("M","").strip()
        a = float(_clean(actual_str))
        t = float(_clean(target_str))
        if t == 0:
            return "black"
        if a >= t:
            return "green"
        elif a >= 0.70 * t:
            return "amber"
        else:
            return "red"
    except Exception:
        return "black"

def _get_rgb(color_str):
    return {"green": GREEN, "amber": AMBER, "red": RED}.get(color_str, BLACK)

def _calc_by_region(df_snap, base_mask, num_mask):
    result = {}
    for reg in REGIONS:
        rm = df_snap[REG_COL] == reg
        den = df_snap[base_mask & rm]["LEADING_END_CUSTOMER_ID"].nunique()
        num = df_snap[num_mask & rm]["LEADING_END_CUSTOMER_ID"].nunique()
        result[reg] = _pct_str(num, den)
    den_all = df_snap[base_mask]["LEADING_END_CUSTOMER_ID"].nunique()
    num_all = df_snap[num_mask]["LEADING_END_CUSTOMER_ID"].nunique()
    result["APAC"] = _pct_str(num_all, den_all)
    return result

def _compute_kpis_for_snapshot(df_snap):
    """Compute all KPI values for a given snapshot DataFrame slice."""
    b26     = df_snap["BASELINE_CUSTOMER_2026"].astype(str).str.strip() == "True"
    b_all   = df_snap["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip().isin(["True","ELIG_nCOV"])
    b_covd  = df_snap["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip() == "True"
    prem    = df_snap["PREMIUM_AI_CUSTOMER"].astype(str).str.strip().str.upper() == "YES"
    sig_lic = df_snap["SIGNAVIO_LICENSED"].astype(str).str.strip().str.upper() == "YES"
    lix_lic = df_snap["LEANIX_LICENSED"].astype(str).str.strip().str.upper() == "YES"

    vals = {}
    vals["joule"]    = _calc_by_region(df_snap, b26 & prem,    b26 & prem    & (df_snap["AI_ACTIVE_JOULE"].astype(str).str.upper() == "YES"))
    vals["ai_plan"]  = _calc_by_region(df_snap, b26,           b26 & df_snap["AI_ADOPTION_PLAN_ACTIVATED"].apply(_is_true))
    vals["rwsm"]     = _calc_by_region(df_snap, b26,           b26 & df_snap["RWSM_DASHBOARD_ACTIVATED"].apply(_is_true))
    vals["risk"]     = _calc_by_region(df_snap, b26,           b26 & df_snap["HAS_PCE_RISK_ASSESSMENT"].apply(_is_true))
    vals["calm"]     = _calc_by_region(df_snap, b26,           b26 & df_snap["CALM_TENANT_ACTIVATED"].apply(_is_true))
    vals["gainsight"]= dict(FROZEN["gainsight"])  # frozen
    vals["signavio"] = _calc_by_region(df_snap, b26 & sig_lic, b26 & sig_lic & df_snap["SIGNAVIO_GT_5_PERCENT_CRATIO"].apply(_is_true))
    vals["leanix"]   = _calc_by_region(df_snap, b26 & lix_lic, b26 & lix_lic & df_snap["LEAN_IX_USAGE"].apply(_is_true))

    # Coverage
    cov = {}
    for reg in ALL_REGIONS:
        if reg == "APAC":
            elig = df_snap[b_all]; covd = df_snap[b_covd]
        else:
            rm = df_snap[REG_COL] == reg
            elig = df_snap[b_all & rm]; covd = df_snap[b_covd & rm]
        tc = elig["LEADING_END_CUSTOMER_ID"].nunique()
        ec = covd["LEADING_END_CUSTOMER_ID"].nunique()
        ta = elig["EXIT_ACV"].sum()
        ea = covd["EXIT_ACV"].sum()
        ne = covd[covd["EA_NAME"].apply(lambda x: str(x).strip() not in ("#","","nan"))]["EA_NAME"].nunique()
        cov[reg] = {
            "kpi9":  _pct_str(ec, tc),
            "kpi10": _pct_str(ea, ta),
            "kpi11": f"{round(ec/ne,2)}" if ne else "N/A",
            "kpi12": f"{round(ea/ne/1e6,1)} M€" if ne else "N/A",
        }
    vals["ea_cust"]   = {r: cov[r]["kpi9"]  for r in ALL_REGIONS}
    vals["ea_acv"]    = {r: cov[r]["kpi10"] for r in ALL_REGIONS}
    vals["ea_density"]= {r: cov[r]["kpi11"] for r in ALL_REGIONS}
    vals["ea_acv_den"]= {r: cov[r]["kpi12"] for r in ALL_REGIONS}

    vals["leanix_enable"] = dict(FROZEN["leanix_enable"])
    vals["ea_learning"]   = dict(FROZEN["ea_learning"])
    return vals


# ── Main compute function ─────────────────────────────────────────────────────

def compute_kpis(excel_path):
    """Load Impact KPI.xlsx and return full KPI dict + timestamp."""
    df = pd.read_excel(excel_path, sheet_name="Impact KPI")
    ts_col = pd.to_datetime(df["TIME_STAMP"], errors="coerce")

    # Sort all unique timestamps; latest = most recent snapshot
    ts_sorted = sorted(ts_col.dropna().unique())
    latest_ts = ts_sorted[-1]
    latest_pd = pd.Timestamp(latest_ts)

    # Prior month = FIRST snapshot in the calendar month immediately before the latest snapshot's month
    # e.g. latest=2026-08-07 → prior_month=(2026,7) → use 2026-07-01 (first July snapshot)
    prev_month_year  = latest_pd.year if latest_pd.month > 1 else latest_pd.year - 1
    prev_month_month = latest_pd.month - 1 if latest_pd.month > 1 else 12
    prior_ts = None
    for ts in ts_sorted:
        t = pd.Timestamp(ts)
        if (t.year, t.month) == (prev_month_year, prev_month_month):
            prior_ts = ts  # first hit in prior month = earliest snapshot that month
            break

    df_l = df[ts_col == latest_ts].copy()
    timestamp = pd.Timestamp(latest_ts).strftime("%d/%m/%Y")

    vals = _compute_kpis_for_snapshot(df_l)

    # Compute prior-month values dynamically from the second-most-recent snapshot
    if prior_ts is not None:
        df_p = df[ts_col == prior_ts].copy()
        prior_vals = _compute_kpis_for_snapshot(df_p)
    else:
        prior_vals = vals  # no prior: delta will be zero

    # Build delta rows
    deltas = {}
    for key in vals:
        deltas[key] = {r: _delta(vals[key].get(r,""), prior_vals[key].get(r,"")) for r in ALL_REGIONS}

    # Store prior vals and prior timestamp for MoM bracket in slide 2
    prior_ts_str = pd.Timestamp(prior_ts).strftime("%d/%m/%Y") if prior_ts is not None else ""

    return {
        "timestamp":  timestamp,
        "prior_ts":   prior_ts_str,
        "values":     vals,
        "deltas":     deltas,
        "prior_vals": prior_vals,
        "targets":    Q3_TARGETS,
    }


def _kpi_row_json(key, label, kpi_data):
    vals    = kpi_data["values"][key]
    deltas  = kpi_data["deltas"][key]
    target  = kpi_data["targets"][key]
    row = {"kpi": label, "target": target, "values": {}, "deltas": {}, "colors": {}}
    for r in ALL_REGIONS:
        v = vals.get(r, "")
        row["values"][r]  = v
        row["deltas"][r]  = deltas.get(r, "(⬤)")
        row["colors"][r]  = get_color_str(v, target)
    return row


def kpi_data_to_json(kpi_data):
    """Convert compute_kpis() output to the JSON structure the dashboard expects.

    KPI order (dashboard + PPT):
      Joule → Signavio → LeanIX → Coverage (# + ACV) → Density (cust + ACV) → Gainsight → Learning
    """
    mk = _kpi_row_json

    # Flat ordered list — each entry has an extra 'category' key for grouping display
    SLIDE1_ORDER = [
        # Business Impact (4 KPIs)
        ("joule",        "Joule Activated",              "Business Impact"),
        ("ai_plan",      "AI Adoption Plan",             "Business Impact"),
        ("rwsm",         "RwSM Dashboard activated",     "Business Impact"),
        ("risk",         "Risk Assessment available",    "Business Impact"),
        # Tool Adoption (4 KPIs)
        ("calm",         "CALM tenant activated",        "Tool Adoption"),
        ("gainsight",    "Gainsight Usage",              "Tool Adoption"),
        ("signavio",     "Active Signavio Usage",        "Tool Adoption"),
        ("leanix",       "Active LeanIX usage",          "Tool Adoption"),
        # Coverage (4 KPIs)
        ("ea_cust",      "EA Coverage (# of Customers)", "Coverage"),
        ("ea_acv",       "EA Coverage (Exit ACV)",       "Coverage"),
        ("ea_density",   "EA Customer Density",          "Coverage"),
        ("ea_acv_den",   "EA ACV Density",               "Coverage"),
        # Learning (2 KPIs)
        ("leanix_enable","Lean IX Enablement",           "Learning"),
        ("ea_learning",  "EA Learning Journey",          "Learning"),
    ]

    # Build rows with category tag embedded
    slide1_rows = []
    for key, label, cat in SLIDE1_ORDER:
        row = mk(key, label, kpi_data)
        row["category"] = cat
        row["key"] = key
        row["frozen"] = key in FROZEN
        slide1_rows.append(row)

    # Also keep the old dict format for PPT builder (by category)
    slide1 = {
        "Business Impact": [r for r in slide1_rows if r["category"] == "Business Impact"],
        "Tool Adoption":   [r for r in slide1_rows if r["category"] == "Tool Adoption"],
        "Coverage":        [r for r in slide1_rows if r["category"] == "Coverage"],
        "Learning":        [r for r in slide1_rows if r["category"] == "Learning"],
        "_ordered": slide1_rows,   # flat ordered list for the dashboard
    }

    v  = kpi_data["values"]
    pv = kpi_data.get("prior_vals", v)
    t  = kpi_data["targets"]

    def _mom_val(key):
        """Return MoM bracket string for APAC, or '' if zero/unavailable."""
        return _mom_bracket(v[key]["APAC"], pv[key]["APAC"])

    def _commentary(key, extra_text):
        apac_val  = v[key]["APAC"]
        prior_val = pv[key]["APAC"]
        bracket   = _mom_bracket(apac_val, prior_val)
        mom_str   = f" ({bracket})" if bracket else ""
        tgt_val   = t[key]
        base = f"APAC YTD attainment of {apac_val}{mom_str} vs Q3 target of {tgt_val}."
        return f"{base} {extra_text}".strip()

    FROZEN_COMMENTARY = "Frozen — no current target set."

    # slide2 rows — same order as SLIDE1_ORDER (all 14 KPIs)
    slide2 = [
        # Business Impact
        {"category":"Business Impact","kpi":"Joule Activated",
         "status": get_color_str(v["joule"]["APAC"], t["joule"]),
         "mom": _mom_val("joule"), "frozen": False,
         "commentary": _commentary("joule", "Will progressively achieve the RBT targets. Lack of clarity on SAP managed Joule and shift from premium to Base hindering progress.")},

        {"category":"","kpi":"AI Adoption Plan",
         "status": get_color_str(v["ai_plan"]["APAC"], t["ai_plan"]),
         "mom": _mom_val("ai_plan"), "frozen": False,
         "commentary": _commentary("ai_plan", "Will progressively achieve the RBT targets. Should become a part of EA MXP.")},

        {"category":"","kpi":"RwSM Dashboard activated",
         "status": get_color_str(v["rwsm"]["APAC"], t["rwsm"]),
         "mom": _mom_val("rwsm"), "frozen": False,
         "commentary": _commentary("rwsm", "Will progressively achieve the RBT targets. Challenge on availability of data has been overcome.")},

        {"category":"","kpi":"Risk Assessment available",
         "status": get_color_str(v["risk"]["APAC"], t["risk"]),
         "mom": _mom_val("risk"), "frozen": False,
         "commentary": _commentary("risk", "")},

        # Tool Adoption
        {"category":"Tool Adoption","kpi":"CALM tenant activated",
         "status": get_color_str(v["calm"]["APAC"], t["calm"]),
         "mom": _mom_val("calm"), "frozen": False,
         "commentary": _commentary("calm", "On track.")},

        {"category":"","kpi":"Gainsight Usage",
         "status": "frozen", "mom": "", "frozen": True,
         "commentary": FROZEN_COMMENTARY},

        {"category":"","kpi":"Active Signavio Usage",
         "status": get_color_str(v["signavio"]["APAC"], t["signavio"]),
         "mom": _mom_val("signavio"), "frozen": False,
         "commentary": _commentary("signavio", "Focus on the customer instances for the next quarter.")},

        {"category":"","kpi":"Active LeanIX usage",
         "status": get_color_str(v["leanix"]["APAC"], t["leanix"]),
         "mom": _mom_val("leanix"), "frozen": False,
         "commentary": _commentary("leanix", "Focus on the customer instances for the next quarter.")},

        # Coverage
        {"category":"Coverage","kpi":"EA Coverage (# of Customers)",
         "status": get_color_str(v["ea_cust"]["APAC"], t["ea_cust"]),
         "mom": _mom_val("ea_cust"), "frozen": False,
         "commentary": _commentary("ea_cust", "Coverage calculated based on eligible customer baseline.")},

        {"category":"","kpi":"EA Coverage (Exit ACV)",
         "status": get_color_str(v["ea_acv"]["APAC"], t["ea_acv"]),
         "mom": _mom_val("ea_acv"), "frozen": False,
         "commentary": _commentary("ea_acv", "Coverage calculated based on eligible ACV baseline.")},

        {"category":"","kpi":"EA Customer Density",
         "status": get_color_str(v["ea_density"]["APAC"], t["ea_density"]),
         "mom": _mom_val("ea_density"), "frozen": False,
         "commentary": _commentary("ea_density", "Density compares YTD customer density vs target.")},

        {"category":"","kpi":"EA ACV Density",
         "status": get_color_str(v["ea_acv_den"]["APAC"], t["ea_acv_den"]),
         "mom": _mom_val("ea_acv_den"), "frozen": False,
         "commentary": _commentary("ea_acv_den", "Density compares YTD ACV density vs target.")},

        # Learning
        {"category":"Learning","kpi":"Lean IX Enablement",
         "status": "frozen", "mom": "", "frozen": True,
         "commentary": FROZEN_COMMENTARY},

        {"category":"","kpi":"EA Learning Journey",
         "status": "frozen", "mom": "", "frozen": True,
         "commentary": FROZEN_COMMENTARY},

        {"category":"Global Observations*","kpi":"","status":"black","mom":"","frozen":False,
         "commentary":"Several KPIs below target are in line with expected Q3 progress. EA assignments finalised, KPI calculation updated. Coverage based on number of accounts shows low % due to broad eligible Business Partners. Coverage on ExitACV shows focus on most important targets is effective."},
    ]

    return {"timestamp": kpi_data["timestamp"], "prior_ts": kpi_data.get("prior_ts",""), "slide1": slide1, "slide2": slide2}


# ── PPTX builder ─────────────────────────────────────────────────────────────

def _remove_fill(cell):
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is not None:
        for tag in [qn("a:solidFill"),qn("a:gradFill"),qn("a:noFill"),qn("a:blipFill"),qn("a:pattFill")]:
            for el in tcPr.findall(tag):
                tcPr.remove(el)

def _set_text(cell, text):
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.add_run().text = text

def _set_font_color(cell, rgb):
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb

def _find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def build_pptx(kpi_data, template_path, output_path):
    """Write updated PPTX from kpi_data dict to output_path."""
    prs = Presentation(template_path)
    vals   = kpi_data["values"]
    deltas = kpi_data["deltas"]
    ts     = kpi_data["timestamp"]

    COL_MAP = {"ANZ":2,"G.China":3,"India":4,"Japan":5,"Korea":6,"South East Asia":7,"APAC":8}

    # ── Slide 1 ──
    s1 = prs.slides[0]
    ts_shape = _find_shape(s1, "Table 47")
    if ts_shape and ts_shape.has_table:
        _set_text(ts_shape.table.cell(0,0), f"Last Updated: {ts}")

    tbl = _find_shape(s1, "Table 16").table

    KPI_ROWS = [
        (2,  3,  "joule"),
        (4,  5,  "ai_plan"),
        (6,  7,  "rwsm"),
        (8,  9,  "risk"),
        (11, 12, "calm"),
        (13, 14, "gainsight"),
        (15, 16, "signavio"),
        (17, 18, "leanix"),
        (20, 21, "ea_cust"),
        (22, 23, "ea_acv"),
        (24, 25, "ea_density"),
        (26, 27, "ea_acv_den"),
        (29, 30, "leanix_enable"),
        (31, 32, "ea_learning"),
    ]

    for (data_row, delta_row, key) in KPI_ROWS:
        tgt = Q3_TARGETS[key]
        _set_text(tbl.cell(data_row, 9), tgt)
        for reg, col in COL_MAP.items():
            v = vals[key].get(reg, "")
            d = deltas[key].get(reg, "(⬤)")
            clr = _get_rgb(get_color_str(v, tgt))
            for row, txt in [(data_row, v), (delta_row, d)]:
                cell = tbl.cell(row, col)
                _remove_fill(cell)
                _set_text(cell, txt)
                _set_font_color(cell, clr)

    # ── Slide 2 ──
    s2 = prs.slides[1]
    ts2 = _find_shape(s2, "Table 47")
    if ts2 and ts2.has_table:
        _set_text(ts2.table.cell(0,0), f"Last Updated: {ts}")

    qtbl = _find_shape(s2, "Table 7").table
    jd = kpi_data_to_json(kpi_data)

    for i, row_data in enumerate(jd["slide2"][:-1]):  # skip global obs
        row_idx = i + 1
        try:
            # Update commentary (col 3)
            cell_comm = qtbl.cell(row_idx, 3)
            tf = cell_comm.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = row_data["commentary"]
            elif tf.paragraphs:
                tf.paragraphs[0].add_run().text = row_data["commentary"]

            # Update traffic light dot color (col 2) using same green/amber/red logic
            cell_dot = qtbl.cell(row_idx, 2)
            dot_rgb  = _get_rgb(row_data["status"])
            _set_font_color(cell_dot, dot_rgb)
        except Exception:
            pass

    prs.save(output_path)


# ── KPI definitions for drilldown ────────────────────────────────────────────
# Each KPI: base_mask_fn(df) → boolean Series, num_mask_fn(df) → boolean Series, status_col, status_label
KPI_DRILLDOWN_DEFS = {
    "joule":     {"base": lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["PREMIUM_AI_CUSTOMER"].astype(str).str.upper().str.strip()=="YES"),
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["PREMIUM_AI_CUSTOMER"].astype(str).str.upper().str.strip()=="YES") & (d["AI_ACTIVE_JOULE"].astype(str).str.upper().str.strip()=="YES"),
                  "status_col": "AI_ACTIVE_JOULE", "status_label": "Joule Active"},
    "ai_plan":   {"base": lambda d: d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True",
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["AI_ADOPTION_PLAN_ACTIVATED"].astype(str).str.strip()=="True"),
                  "status_col": "AI_ADOPTION_PLAN_ACTIVATED", "status_label": "AI Plan Active"},
    "rwsm":      {"base": lambda d: d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True",
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["RWSM_DASHBOARD_ACTIVATED"].astype(str).str.strip()=="True"),
                  "status_col": "RWSM_DASHBOARD_ACTIVATED", "status_label": "RwSM Active"},
    "risk":      {"base": lambda d: d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True",
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["HAS_PCE_RISK_ASSESSMENT"].astype(str).str.strip()=="True"),
                  "status_col": "HAS_PCE_RISK_ASSESSMENT", "status_label": "Risk Assessment"},
    "calm":      {"base": lambda d: d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True",
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["CALM_TENANT_ACTIVATED"].astype(str).str.strip()=="True"),
                  "status_col": "CALM_TENANT_ACTIVATED", "status_label": "CALM Active"},
    "signavio":  {"base": lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["SIGNAVIO_LICENSED"].astype(str).str.upper().str.strip()=="YES"),
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["SIGNAVIO_LICENSED"].astype(str).str.upper().str.strip()=="YES") & (d["SIGNAVIO_GT_5_PERCENT_CRATIO"].astype(str).str.strip()=="True"),
                  "status_col": "SIGNAVIO_GT_5_PERCENT_CRATIO", "status_label": "Active Usage"},
    "leanix":    {"base": lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["LEANIX_LICENSED"].astype(str).str.upper().str.strip()=="YES"),
                  "num":  lambda d: (d["BASELINE_CUSTOMER_2026"].astype(str).str.strip()=="True") & (d["LEANIX_LICENSED"].astype(str).str.upper().str.strip()=="YES") & (d["LEAN_IX_USAGE"].astype(str).str.strip()=="True"),
                  "status_col": "LEAN_IX_USAGE", "status_label": "Active Usage"},
    "ea_cust":   {"base": lambda d: d["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip().isin(["True","ELIG_nCOV"]),
                  "num":  lambda d: d["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip()=="True",
                  "status_col": "BASELINE_CUSTOMER_DYNAMIC", "status_label": "EA Covered"},
    "ea_acv":    {"base": lambda d: d["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip().isin(["True","ELIG_nCOV"]),
                  "num":  lambda d: d["BASELINE_CUSTOMER_DYNAMIC"].astype(str).str.strip()=="True",
                  "status_col": "BASELINE_CUSTOMER_DYNAMIC", "status_label": "EA Covered"},
}


def drilldown(excel_path, kpi_key, region, snapshot_date=None, df_cached=None):
    """Return account-level rows for a given KPI + region + snapshot.
    Pass df_cached to avoid re-reading the Excel file."""
    if df_cached is not None:
        df = df_cached
    else:
        df = pd.read_excel(excel_path, sheet_name="Impact KPI")
    ts_col = pd.to_datetime(df["TIME_STAMP"], errors="coerce")

    if snapshot_date:
        target_ts = pd.Timestamp(snapshot_date)
        df_snap = df[ts_col == target_ts].copy()
    else:
        latest_ts = ts_col.max()
        df_snap = df[ts_col == latest_ts].copy()

    if kpi_key not in KPI_DRILLDOWN_DEFS:
        return {"error": f"No drilldown for {kpi_key}"}

    defn   = KPI_DRILLDOWN_DEFS[kpi_key]
    base_m = defn["base"](df_snap)
    num_m  = defn["num"](df_snap)

    # Filter by region
    if region and region != "APAC":
        reg_m  = df_snap[REG_COL] == region
        base_m = base_m & reg_m
        num_m  = num_m  & reg_m

    df_base = df_snap[base_m].copy()
    df_base["_activated"] = num_m[base_m]

    # Include account name + id, region, EA, baseline flags, KPI status, ACV
    keep_cols = ["LEADING_END_CUSTOMER", "LEADING_END_CUSTOMER_ID", REG_COL, "EA_NAME",
                 "BASELINE_CUSTOMER_2026", "BASELINE_CUSTOMER_DYNAMIC",
                 defn["status_col"], "EXIT_ACV"]
    keep_cols = [c for c in keep_cols if c in df_base.columns]
    df_out = df_base[keep_cols + ["_activated"]].drop_duplicates(subset=["LEADING_END_CUSTOMER_ID"])
    df_out = df_out.sort_values("_activated", ascending=False)

    # Rename for display
    df_out = df_out.rename(columns={
        "LEADING_END_CUSTOMER":    "Account Name",
        "LEADING_END_CUSTOMER_ID": "Account ID",
        REG_COL:                   "Region",
        "EA_NAME":                 "EA Name",
        "BASELINE_CUSTOMER_2026":  "B2026",
        "BASELINE_CUSTOMER_DYNAMIC": "B Dynamic",
        defn["status_col"]:        defn["status_label"],
        "EXIT_ACV":                "Exit ACV (€)",
        "_activated":              "KPI Met",
    })

    # Format Exit ACV as readable number
    if "Exit ACV (€)" in df_out.columns:
        df_out["Exit ACV (€)"] = df_out["Exit ACV (€)"].apply(
            lambda x: f"{float(x)/1e6:.1f}M" if str(x).replace('.','').replace('-','').isdigit() or _is_numeric(x) else x
        )

    total   = len(df_out)
    met     = int(df_out["KPI Met"].sum())
    records = df_out.fillna("").astype(str).to_dict(orient="records")

    return {
        "kpi":     kpi_key,
        "region":  region,
        "snapshot": str(pd.Timestamp(df_snap["TIME_STAMP"].iloc[0]).date()) if len(df_snap) else "",
        "total":   total,
        "met":     met,
        "pct":     round(100*met/total) if total else 0,
        "columns": list(df_out.columns),
        "rows":    records,
    }

def _is_numeric(x):
    try: float(x); return True
    except: return False


def trend(excel_path, kpi_key, df_cached=None):
    """Return all-region % for a KPI across all snapshots. Uses df_cached if provided."""
    df = df_cached if df_cached is not None else pd.read_excel(excel_path, sheet_name="Impact KPI")
    ts_col = pd.to_datetime(df["TIME_STAMP"], errors="coerce")
    ts_sorted = sorted(ts_col.dropna().unique())

    points = []
    for ts in ts_sorted:
        df_snap = df[ts_col == ts].copy()
        vals    = _compute_kpis_for_snapshot(df_snap)
        apac    = vals.get(kpi_key, {}).get("APAC", "N/A")
        regions = {r: vals.get(kpi_key, {}).get(r, "N/A") for r in REGIONS}
        points.append({
            "date":   pd.Timestamp(ts).strftime("%d/%m"),
            "APAC":   apac,
            **regions,
        })

    return {
        "kpi":     kpi_key,
        "target":  Q3_TARGETS.get(kpi_key, ""),
        "points":  points,
        "regions": ALL_REGIONS,
    }

